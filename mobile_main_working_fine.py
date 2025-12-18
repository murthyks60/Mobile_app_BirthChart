"""
consolidated_birthchart_generator.py
----------------------------------------------
Goal: Collect global birth data (POB anywhere in the world), compute sidereal basics
(Julian Day, Ayanamsa, planetary positions), Panchanga, and now Nakshatra–pāda transitions.
Populate into base_dict and feed to PPTX generator.
"""

import pandas as pd
import swisseph as swe
from datetime import datetime
from datetime import timedelta
from zoneinfo import ZoneInfo
from datetime import timezone
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
import math
from math import *
import requests
from timezonefinder import TimezoneFinder
import pytz
import datetime as dt




#ist_tz = timezone(timedelta(hours=5, minutes=30))


SHOW_FINAL_OUTPUT = True   # set False to suppress even the final print

# --- Panchanga constants ---
TITHI_NAMES = [
    "Pratipada", "Dwitiya", "Tritiya", "Chaturthi", "Panchami", "Shashti",
    "Saptami", "Ashtami", "Navami", "Dashami", "Ekadashi", "Dwadashi",
    "Trayodashi", "Chaturdashi", "Purnima",
    "Pratipada (Krishna)", "Dwitiya (Krishna)", "Tritiya (Krishna)", "Chaturthi (Krishna)",
    "Panchami (Krishna)", "Shashti (Krishna)", "Saptami (Krishna)", "Ashtami (Krishna)",
    "Navami (Krishna)", "Dashami (Krishna)", "Ekadashi (Krishna)", "Dwadashi (Krishna)",
    "Trayodashi (Krishna)", "Chaturdashi (Krishna)", "Amavasya"
]

NAKSHATRA_NAMES = [
    "Ashwini","Bharani","Krittika","Rohini","Mrigashira","Ardra","Punarvasu","Pushya",
    "Ashlesha","Magha","Purva Phalguni","Uttara Phalguni","Hasta","Chitra","Swati","Vishakha",
    "Anuradha","Jyeshtha","Mula","Purva Ashadha","Uttara Ashadha","Shravana","Dhanishta",
    "Shatabhisha","Purva Bhadrapada","Uttara Bhadrapada","Revati"
]

YOGA_NAMES = [
    "Vishkumbha","Preeti","Ayushman","Saubhagya","Shobhana","Atiganda","Sukarma","Dhriti",
    "Shoola","Ganda","Vriddhi","Dhruva","Vyaghata","Harshana","Vajra","Siddhi","Vyatipata",
    "Variyana","Parigha","Shiva","Siddha","Sadhya","Shubha","Shukla","Brahma","Indra","Vaidhriti"
]

KARANA_NAMES = [
    "Bava","Balava","Kaulava","Taitila","Garaja","Vanija","Vishti",
    "Shakuni","Chatushpada","Nagava","Kimstughna"
]

KARANA_SEQUENCE = (
    ["Kimstughna"]
    + ["Bava", "Balava", "Kaulava", "Taitila", "Garaja", "Vanija", "Vishti"] * 8
    + ["Shakuni", "Chatushpada", "Naga"]
)
planet_map = {
                "Sun": "Su", "Moon": "Mo", "Mercury": "Me", "Venus": "Ve",
                "Mars": "Ma", "Jupiter": "Ju", "Saturn": "Sa",
                "Rahu": "Ra", "Ketu": "Ke", "Ascendant": "Asc"
}

WEEKDAY_NAMES = ["Monday","Tuesday","Wednesday","Thursday","Friday","Saturday","Sunday"]

SIGNS = ["Aries","Taurus","Gemini","Cancer","Leo","Virgo",
         "Libra","Scorpio","Sagittarius","Capricorn","Aquarius","Pisces"]

# Set ephemeris path
swe.set_ephe_path("C:/Users/murth/Downloads/Swisseph_extractions_13Nov25/ephe")

# Planets we compute

# Explicit mapping from Swisseph planet constants to our base_dict keys
PLANET_KEYS = {
    swe.SUN:     "SUN",
    swe.MOON:    "MOON",
    swe.MARS:    "MAR",
    swe.MERCURY: "MER",
    swe.JUPITER: "JUP",
    swe.VENUS:   "VEN",
    swe.SATURN:  "SAT",
}

TELUGU_YEAR_NAMES = [
    "Prabhava", "Vibhava", "Śukla", "Pramōdyuta", "Prajōtpatti", "Āṅgīrasa", "Śrīmukha", "Bhava",
    "Yuva", "Dhāta", "Īśvara", "Bahudhānya", "Pramādhi", "Vikrama", "Vr̥ṣa", "Citrabhānu",
    "Svabhānu", "Tāraṇa", "Pārthiva", "Vyaya", "Sarvajittu", "Sarvadhāri", "Virōdhi", "Vikr̥ti",
    "Khara", "Nandana", "Vijaya", "Jaya", "Manmadha", "Durmukhi", "Hēvaḷambi", "Viḷambi",
    "Vikāri", "Śārvari", "Plava", "Śubhakr̥ttu", "Śōbhakr̥ttu", "Krōdhi", "Viśvāvasu", "Parābhava",
    "Plavaṅga", "Kīlaka", "Saumya", "Sādhāraṇa", "Virōdhikr̥ttu", "Paridhāvi", "Pramādīca", "Ānanda",
    "Rākṣasa", "Nala", "Piṅgaḷa", "Kāḷayukti", "Siddhārthi", "Raudri", "Durmati", "Dundubhi",
    "Rudhirōdgāri", "Raktākṣi", "Krōdhana", "Akṣaya"
]
PLANETS = {
        "Sun": swe.SUN,
        "Moon": swe.MOON,
        "Mars": swe.MARS,
        "Mercury": swe.MERCURY,
        "Jupiter": swe.JUPITER,
        "Venus": swe.VENUS,
        "Saturn": swe.SATURN,
        "Rahu": swe.MEAN_NODE,   # North Node
        "Ketu": swe.TRUE_NODE    # South Node (can also compute as 180° opposite Rahu)
}

swe.set_sid_mode(swe.SIDM_LAHIRI, 0, 0)
flags = swe.FLG_SWIEPH | swe.FLG_SPEED | swe.FLG_SIDEREAL

# --- Nakshatra pāda transitions helpers ---
NAK = 360.0 / 27.0
PADA = NAK / 4.0

ZODIAC = ["Ar","Ta","Ge","Cn","Le","Vi","Li","Sc","Sg","Cp","Aq","Pi"]

# Example TITHI_NAMES list (index 0–29)
TITHI_NAMES = [
    "Shukla Prathama", "Shukla Dwitiya", "Shukla Tritiya", "Shukla Chaturthi", "Shukla Panchami",
    "Shukla Shashthi", "Shukla Saptami", "Shukla Ashtami", "Shukla Navami", "Shukla Dashami",
    "Shukla Ekadashi", "Shukla Dwadashi", "Shukla Trayodashi", "Shukla Chaturdashi", "Purnima",
    "Krishna Prathama", "Krishna Dwitiya", "Krishna Tritiya", "Krishna Chaturthi", "Krishna Panchami",
    "Krishna Shashthi", "Krishna Saptami", "Krishna Ashtami", "Krishna Navami", "Krishna Dashami",
    "Krishna Ekadashi", "Krishna Dwadashi", "Krishna Trayodashi", "Krishna Chaturdashi", "Amavasya"
]

def get_sunrise_sunset(city_name, dt_obj):
    """
    Compute local sunrise and sunset times for a given city and datetime.
    Uses lookup_city() and tz_to_offset() helpers.
    """

    # Step 1: Lookup city details
    lat, lon, utc_offset, tz_string = lookup_city(city_name)

    # Step 2: Convert datetime to Julian Day
    jd_ut = swe.julday(
        dt_obj.year,
        dt_obj.month,
        dt_obj.day,
        dt_obj.hour + dt_obj.minute / 60.0 + dt_obj.second / 3600.0,
        swe.GREG_CAL
    )

    # Step 3: Compute UTC sunrise and sunset
    sunrise_result = swe.rise_trans(jd_ut, swe.SUN, lon, lat, rsmi=swe.CALC_RISE)
    sunset_result  = swe.rise_trans(jd_ut, swe.SUN, lon, lat, rsmi=swe.CALC_SET)

    sunrise_utc = sunrise_result[1][0]
    sunset_utc  = sunset_result[1][0]

    # Step 4: Apply timezone offset
    offset = tz_to_offset(
        tz_string,
        dt_obj.strftime("%Y-%m-%d"),
        dt_obj.strftime("%H:%M:%S")
    )

    sunrise_local = sunrise_utc + timedelta(hours=offset)
    sunset_local  = sunset_utc + timedelta(hours=offset)

    return sunrise_local, sunset_local

def get_tithi_transitions(year, month, day, hour, minute, second, utc_offset, steps=24):
    """Compute Tithi transitions for the given day."""
    jd_start = swe.julday(year, month, day, hour + minute/60 + second/3600)
    transitions = []
    for i in range(steps):
        jd = jd_start + i/24.0  # step in hours
        moon_lon = swe.calc_ut(jd, swe.MOON)[0][0]
        sun_lon  = swe.calc_ut(jd, swe.SUN)[0][0]
        ayanamsa = swe.get_ayanamsa(jd)
        moon_sid = (moon_lon - ayanamsa) % 360
        sun_sid  = (sun_lon - ayanamsa) % 360
        sep = (moon_sid - sun_sid) % 360
        tithi_index = int(sep // 12) % len(TITHI_NAMES)
        tithi_name  = TITHI_NAMES[tithi_index]
        dt = swe.revjul(jd, swe.GREG_CAL)
        dt_obj = datetime(dt[0], dt[1], dt[2], int((dt[3])%24), int((dt[3]%1)*60))
        transitions.append({"tithi": tithi_name, "datetime": dt_obj})
    return transitions

def compute_tithi(moon_sidereal, sun_sidereal, year, month, day, hour, minute, second, utc_offset):
    # Current tithi name
    tithi_angle = (moon_sidereal - sun_sidereal) % 360
    tithi_index = int(tithi_angle // 12)
    tithi_name = TITHI_NAMES[tithi_index]

    # Use transitions to find end time and next tithi
    transitions = get_tithi_transitions(year, month, day, hour, minute, second, utc_offset)
    # Find the first transition where tithi changes
    current = None
    next_tithi_name = None
    tithi_end_time = None
    for i in range(1, len(transitions)):
        if transitions[i]["tithi"] != transitions[i-1]["tithi"]:
            current = transitions[i-1]["tithi"]
            next_tithi_name = transitions[i]["tithi"]
            tithi_end_time = transitions[i]["datetime"]
            break

    # Build message
    if tithi_end_time and next_tithi_name:
        tithi_end_message = f"{tithi_name} ends at {tithi_end_time.strftime('%Y-%m-%d %H:%M')}, then {next_tithi_name} begins"
    else:
        tithi_end_message = f"{tithi_name} continues all day"

    return tithi_name, tithi_end_time, next_tithi_name, tithi_end_message

def tz_to_offset(tz_string, dob, tob):
    year, month, day = map(int, dob.split("-"))
    hour, minute, second = map(int, tob.split(":"))
    dt = datetime.datetime(year, month, day, hour, minute, second, tzinfo=ZoneInfo(tz_string))
    offset = dt.utcoffset().total_seconds() / 3600.0
    return offset

def get_all_positions(jd_ut, lat, lon, flags):
    """
    Compute sidereal positions for all planets + Ascendant.
    Returns a list of dicts with keys: name, lon, abs_lon, sign.
    """
    positions = []

    # Planets
    for planet, key in zip(PLANETS, PLANET_KEYS):
        lon, lat_p, dist, speed = swe.calc_ut(jd_ut, planet, flags)
        abs_lon = lon % 360
        sign = SIGNS[int(abs_lon // 30)]
        positions.append({
            "name": key,
            "lon": lon,
            "abs_lon": abs_lon,
            "sign": sign
        })

    # Ascendant (Lagna)
    houses = swe.houses(jd_ut, lat, lon, b'A')  # Placidus system by default
    asc_lon = houses[1][0] % 360                # Ascendant longitude
    asc_sign = SIGNS[int(asc_lon // 30)]
    positions.append({
        "name": "Ascendant",
        "lon": asc_lon,
        "abs_lon": asc_lon,
        "sign": asc_sign
    })

    return positions

def get_sidereal_positions(jd_ut, flags):
    
    positions = {}

    for name, planet in PLANETS.items():
        planet_pos, ret = swe.calc_ut(jd_ut, swe.MOON, flags)
        abs_lon = planet_pos[0] % 360   # Moon’s normalized longitude
        ayanamsa = swe.get_ayanamsa(jd_ut)
        lon = (planet_pos[0] - ayanamsa) % 360
        positions[name] = lon   # assign directly to dict
    return positions

def compute_julian_day(dob, tob, utc_offset):
    year, month, day = map(int, dob.split("-"))
    hour, minute, second = map(int, tob.split(":"))

    # ✅ Force conversion here
    utc_offset = float(utc_offset)

    # Local → UTC
    ut_hour = hour - utc_offset
    ut = ut_hour + minute/60.0 + second/3600.0

    jd_ut = swe.julday(year, month, day, ut, swe.GREG_CAL)
    return jd_ut

def get_telugu_year(gregorian_year):
    base_year = 1927  # Prabhava
    index = (gregorian_year - base_year) % 60
    return TELUGU_YEAR_NAMES[index]

def compute_nakshatra(abs_lon, jd_ut, utc_offset):
    nak_index = int(abs_lon // (360/27))   # which Nakshatra
    nak_name = NAKSHATRA_NAMES[nak_index]
    pada_idx = int((abs_lon % (360/27)) // (360/108)) + 1
    nak_full = f"{nak_name} ({pada_idx})"

    jd_temp = jd_ut
    nak_end = None
    while True:
        jd_temp += 0.01  # step ~15 minutes
        moon_pos, _ = swe.calc_ut(jd_temp, swe.MOON, flags)
        ayanamsa = swe.get_ayanamsa(jd_temp)
        moon_sid = (moon_pos[0] - ayanamsa) % 360

        new_nak_index = int(moon_sid // (360/27))
        new_pada_idx = int((moon_sid % (360/27)) // (360/108)) + 1

        if new_nak_index != nak_index or new_pada_idx != pada_idx:
            year, month, day, hour_float = swe.revjul(jd_temp, swe.GREG_CAL)
            hour = int(hour_float)
            minute = int((hour_float % 1) * 60)
            nak_end = datetime(year, month, day, hour, minute)
            break

    return nak_name, pada_idx, nak_end

def compute_karana(moon_sidereal, sun_sidereal):
    tithi_angle = (moon_sidereal - sun_sidereal) % 360
    karana_index = int(tithi_angle // 6)

    # KARANA_SEQUENCE = (
    #     ["Bava", "Balava", "Kaulava", "Taitila", "Garaja", "Vanija", "Vishti"] * 8
    #     + ["Shakuni", "Chatushpada", "Naga", "Kimstughna"]
    # )
    
    KARANA_SEQUENCE = (
    ["Kimstughna"]
    + ["Bava", "Balava", "Kaulava", "Taitila", "Garaja", "Vanija", "Vishti"] * 8
    + ["Shakuni", "Chatushpada", "Naga"]
    )

    karana = KARANA_SEQUENCE[karana_index]

    # 🧾 Debug prints
    # print("\n🔍 Karana Debug")
    # print(f"🌙 Moon sidereal: {moon_sidereal:.6f}")
    # print(f"🌞 Sun sidereal: {sun_sidereal:.6f}")
    # print(f"📐 Tithi angle (Moon - Sun): {tithi_angle:.6f}")
    # print(f"🔢 Karana index: {karana_index}")
    # print(f"🪔 Karana resolved: {karana}")

    return karana

def compute_yoga(moon_sidereal, sun_sidereal):
    yoga_index = int((moon_sidereal + sun_sidereal) % 360 // (360 / 27))
    return YOGA_NAMES[yoga_index]

def lookup_city(city_name):
    url = "https://nominatim.openstreetmap.org/search"
    params = {"q": city_name, "format": "json", "limit": 1}
    response = requests.get(url, params=params, headers={"User-Agent": "birth-chart-app"})
    data = response.json()
    if not data:
        raise ValueError("City not found")

    lat = float(data[0]["lat"])
    lon = float(data[0]["lon"])

    tf = TimezoneFinder()
    tz_string = tf.timezone_at(lat=lat, lng=lon)
    if tz_string is None:
        tz_string = tf.closest_timezone_at(lat=lat, lng=lon) or "Etc/GMT"

    tz = pytz.timezone(tz_string)
    utc_offset = tz.utcoffset(dt.datetime.now()).total_seconds() / 3600

    return lat, lon, utc_offset, tz_string

def compute_planets(jd_ut, lat, lon, flags=swe.FLG_SWIEPH | swe.FLG_SIDEREAL):
    positions = []
    for name, planet_code in PLANETS.items():   # name is "Sun", planet_code is swe.SUN
        planet_pos, ret = swe.calc_ut(jd_ut, planet_code, flags)
        lon = planet_pos[0]
        abs_lon = lon % 360
        sign = SIGNS[int(abs_lon // 30)]
        positions.append({
            "name": name,
            "lon": lon,
            "abs_lon": abs_lon,
            "sign": sign
        })

    # Ascendant (Lagna)
    houses = swe.houses(jd_ut, lat, lon, b'A')
    asc_lon = houses[1][0] % 360
    asc_sign = SIGNS[int(asc_lon // 30)]
    positions.append({
        "name": "Ascendant",
        "lon": asc_lon,
        "abs_lon": asc_lon,
        "sign": asc_sign
    })
    return positions

def generate_birth_chart(name, dob, tob, city, lat, lon, utc_offset,jd_ut, sidereal_positions):
    # Convert DOB and TOB into datetime objects
    dob_obj = datetime.strptime(dob, "%Y-%m-%d").date()
    tob_obj = datetime.strptime(tob, "%H:%M:%S").time()
    dt_obj  = datetime.combine(dob_obj, tob_obj)
    t0 = datetime.combine(dob_obj, tob_obj)   # this is your reference datetime
    # Weekday
    weekday = dt_obj.strftime("%A")
    chart_dict = {}

    # Telugu year (placeholder – replace with your actual logic)
    telugu_year = get_telugu_year(dob_obj.year)

    # Panchanga calculations using sidereal values
    sun_sidereal  = sidereal_positions["Sun"]
    moon_sidereal = sidereal_positions["Moon"]

    # Tithi
    tithi_name, tithi_end, _, _ = compute_tithi(
        moon_sidereal, sun_sidereal,
        dob_obj.year, dob_obj.month, dob_obj.day,
        tob_obj.hour, tob_obj.minute, tob_obj.second,
        utc_offset
    )
    

    nak_name, nak_pada, nak_end = compute_nakshatra(moon_sidereal, jd_ut, utc_offset)
    # print("DEBUG Nakshatra raw result:", (nak_name, nak_pada, nak_end))
    nak_full = f"{nak_name} (Pada {nak_pada})"


    #print("sidereal_positions keys :", sidereal_positions.keys())
    # Karana
    karana = compute_karana(moon_sidereal, sun_sidereal)
    # print(f"DEBUG Karana: {karana}")
    # Yoga
    yoga_name = compute_yoga(moon_sidereal, sun_sidereal)
    # print(f"DEBUG Yoga: {yoga_name}")
    kar_yog = current_karana_yoga_end(t0, sidereal_positions)
    #print("DEBUG kar_yog dict:", kar_yog)

    # Convert to IST and format
    karana_end_local = kar_yog["karana_end"].astimezone(tz)
    karana_end_str = karana_end_local.strftime("%d-%b-%Y %I:%M %p IST")

    yoga_end_local = kar_yog["yoga_end"].astimezone(tz)
    #chart_dict["YOGA_END"] = yoga_end_local.strftime("%d-%b-%Y %I:%M %p IST")

    # Names
    #chart_dict["KARANA"] = kar_yog["karana_name"]
    #chart_dict["YOGA"]   = kar_yog["yoga_name"]
    #print("DEBUG KARANA_END after assign KS1:", chart_dict["KARANA_END"])
    #print("DEBUG YOGA_END after assign: KS1 ", chart_dict["YOGA_END"])
    # Yoga
    yoga_end_local = kar_yog["yoga_end"].astimezone(tz)
    yoga_end_str = yoga_end_local.strftime("%d-%b-%Y %I:%M %p IST")
    if yoga_end_local.date() > t0.date():
        yoga_end_str += " (continues to next day)"

    #chart_dict["YOGA"]     = kar_yog["yoga_name"]
    #chart_dict["YOGA_END"] = yoga_end_str       # <-- use directly
    #print("DEBUG KARANA_END after assign KS2:", chart_dict["KARANA_END"])
    #print("DEBUG YOGA_END after assign: KS2", chart_dict["YOGA_END"])


    # chart_dict["RAHU_KAALAM"] = ""   # will be filled later
    # chart_dict["DURMUHURTAM"] = ""   # will be filled later

    # Build chart_dict
    # Instead of reassigning chart_dict = {...}
    chart_dict.update({
        "NAME": name,
        "DATE": dob,
        "TIME": tob,
        "PLACE": city,
        "WEEKDAY": weekday,
        "LAT": lat,
        "LONG": lon,
        "TELUGU_YEAR": telugu_year,
        "TITHI": tithi_name,
        "TITHI_END": tithi_end,
        "NAKSHATRA": nak_full,
        "NAK_END": nak_end,
        "KARANA": karana,
        "KARANA_END": karana_end_str,
        "YOGA": yoga_name,
        "YOGA_END": yoga_end_str,
        "PLANETS": sidereal_positions
    })


    # 🔎 Final debug preview
    print("\n📋 Panchanga Preview")
    print("====================================")
    for key, val in chart_dict.items():
        if key in ("chart_data", "PLANETS"):   # skip both
            continue
        print(f"{key:12s}: {val}")
    print("====================================")
    
    for planet, lon in chart_dict["PLANETS"].items():
        dms_str, abs_lon, sign = get_sign_and_abs(lon)
        suffix = planet_map[planet]   # e.g. "Su", "Mo", "Asc"

        chart_dict[f"LONG_{suffix}"] = dms_str          # Column 2 → D:M:S string
        chart_dict[f"ABS_{suffix}"]  = f"{abs_lon:.4f}" # Column 3 → raw float longitude
        chart_dict[f"SIGN_{suffix}"] = sign             # Column 4 → zodiac sign abbreviation
    print("\nPLANET POSITIONS")
    print("========================================================")
    for planet, suffix in planet_map.items():
        print(f"{planet:10s} | LONG={chart_dict.get(f'LONG_{suffix}', '')} "
            f"| ABS={chart_dict.get(f'ABS_{suffix}', '')} "
            f"| SIGN={chart_dict.get(f'SIGN_{suffix}', '')}")
    print("=========================================================")

    return chart_dict
 
def format_end_time(dt_utc, t0_utc, ist_tz):
    # Ensure both are timezone-aware
    if dt_utc.tzinfo is None:
        dt_utc = dt_utc.replace(tzinfo=timezone.utc)
    if t0_utc.tzinfo is None:
        t0_utc = t0_utc.replace(tzinfo=timezone.utc)

    local = dt_utc.astimezone(ist_tz)
    s = local.strftime("%d-%b-%Y %I:%M %p IST")

    if local.date() > t0_utc.astimezone(ist_tz).date():
        s += " (continues to next day)"
    return s

def populate_pptx(chart_dict, template_file="Template_BChart_A4.pptx"):
    fill_pptx(chart_dict, template_file=template_file)
    return template_file

def current_karana_yoga_end(t0, sidereal_positions):
    """
    Compute current Karana/Yoga names and their end times.
    t0: datetime (UTC, tz-aware)
    sidereal_positions: dict with keys including "Sun" and "Moon"
        Values are sidereal longitudes in degrees.
    """

    # --- Current angles and rates ---
    Ls = sidereal_positions["Sun"]   # Sun longitude
    Lm = sidereal_positions["Moon"]  # Moon longitude

    dLs = 0.9856    # Sun mean daily motion (deg/day)
    dLm = 13.176    # Moon mean daily motion (deg/day)

    Delta = (Lm - Ls) % 360.0        # Moon - Sun
    Y     = (Lm + Ls) % 360.0        # Moon + Sun
    rate_delta = dLm - dLs
    rate_yoga  = dLm + dLs

    # --- Next boundaries ---
    karana_step = 6.0
    yoga_step   = 13.3333333333  # 13°20'
    next_karana = math.ceil(Delta / karana_step) * karana_step
    next_yoga   = math.ceil(Y / yoga_step) * yoga_step

    # --- Initial estimates ---
    d_angle_k = (next_karana - Delta) % 360.0
    d_angle_y = (next_yoga - Y) % 360.0
    dt_k_est  = d_angle_k / max(rate_delta, 1e-6)
    dt_y_est  = d_angle_y / max(rate_yoga, 1e-6)

    # --- Refinement helper ---
    def refine(lo, hi, boundary, mode="karana"):
        for _ in range(12):
            tm = lo + (hi - lo) / 2
            delta_days = (tm - t0).total_seconds() / 86400.0
            Ls_m = Ls + dLs * delta_days
            Lm_m = Lm + dLm * delta_days
            val = (Lm_m - Ls_m) % 360.0 if mode == "karana" else (Lm_m + Ls_m) % 360.0
            if val < boundary:
                lo = tm
            else:
                hi = tm
        return hi.replace(tzinfo=timezone.utc)  # attach tzinfo once here

    # --- Refine both ---
    t_k_end = refine(t0, t0 + timedelta(days=dt_k_est), next_karana, mode="karana")
    t_y_end = refine(t0, t0 + timedelta(days=dt_y_est), next_yoga, mode="yoga")

    # --- Karana name mapping ---
    karana_names_cycle = ["Bava", "Balava", "Kaulava", "Taitila",
                          "Garaja", "Vanija", "Vishti"]
    fixed_start = ["Kimstughna"]
    fixed_end   = ["Sakuni", "Chatushpada", "Nagavamsa"]

    half_tithi_index = int(Delta // 6)  # 0–59
    if half_tithi_index == 0:
        karana_name = fixed_start[0]
    elif half_tithi_index >= 57:
        karana_name = fixed_end[half_tithi_index - 57]
    else:
        karana_name = karana_names_cycle[(half_tithi_index - 1) % 7]

    # --- Yoga name mapping ---
    yoga_names = [
        "Vishkambha", "Priti", "Ayushman", "Saubhagya", "Shobhana",
        "Atiganda", "Sukarma", "Dhriti", "Shoola", "Ganda",
        "Vriddhi", "Dhruva", "Vyaghata", "Harshana", "Vajra",
        "Siddhi", "Vyatipata", "Variyan", "Parigha", "Shiva",
        "Siddha", "Sadhya", "Shubha", "Shukla", "Brahma",
        "Indra", "Vaidhriti"
    ]
    yoga_index = int(Y // yoga_step) % 27
    yoga_name = yoga_names[yoga_index]

    return {
        "karana_name": karana_name,
        "karana_end": t_k_end,
        "yoga_name": yoga_name,
        "yoga_end": t_y_end,
        "karana_boundary": next_karana,
        "yoga_boundary": next_yoga
    }

def output_chart(chart_dict):
    print("\n📋 Panchanga Preview")
    print(f"NAME            {chart_dict['NAME']}")
    print(f"DATE            {chart_dict['DATE']}")
    print(f"TIME            {chart_dict['TIME']}")
    print(f"PLACE           {chart_dict['PLACE']}")
    print(f"WEEKDAY         {chart_dict['WEEKDAY']}")
    print(f"TELUGU_YEAR     {chart_dict['TELUGU_YEAR']}")
    print(f"TITHI           {chart_dict['TITHI']}")
    print(f"TITHI_END       {chart_dict['TITHI_END']}")
    print(f"NAKSHATRA       {chart_dict['NAKSHATRA']}")
    print(f"NAK_END         {chart_dict['NAK_END']}")
    print(f"KARANA          {chart_dict['KARANA']}")
    print(f"YOGA            {chart_dict['YOGA']}")

    # 🪐 Planetary Positions
    print("\n🪐 Planetary Positions")
    signs = ["Aries","Taurus","Gemini","Cancer","Leo","Virgo",
             "Libra","Scorpio","Sagittarius","Capricorn","Aquarius","Pisces"]
    for name, lon in chart_dict["PLANETS"].items():
        sign_index = int(lon // 30)
        sign = signs[sign_index]
        print(f"{name:10} {lon:6.2f} {sign}")

def get_sign_and_abs(lon: float):
    """
    Given a planet's longitude (0–360), return:
      - dms_str: degrees, minutes, seconds within its sign (e.g. '21°14\'32.34"')
      - abs_lon: absolute longitude float (e.g. 171.1234)
      - sign: zodiac sign abbreviation (e.g. 'Vi')
    """
    # Absolute longitude float
    abs_lon = lon

    # Degrees within the sign
    deg = int(lon % 30)
    # Minutes within the degree
    mins = int((lon * 60) % 60)
    # Seconds within the minute (fractional)
    secs = (lon * 3600) % 60

    # Format D:M:S string
    dms_str = f"{deg}°{mins:02d}'{secs:05.2f}\""

    # Zodiac sign abbreviation
    sign_index = int(lon // 30)
    sign = ZODIAC[sign_index]

    return dms_str, abs_lon, sign

def fill_pptx(chart_dict, template_file="Template_BChart_A4.pptx"):
    """
    Fill the PPTX template with chart_dict values.
    Populates:
      - Title: "BIRTH CHART OF {{NAME}}" (center aligned, bold, shadow, uppercase)
      - Table 1: Panchanga details
      - Table 2: Rasi Chakra grid (4x4 South Indian style)
      - Table 3: Planetary positions (Longitude, Abs.Longitude, Zodiac Sign)
    Output filename includes user name + datetime stamp.
    """

    prs = Presentation(template_file)
    
    for slide in prs.slides:
        for shape in slide.shapes:

            # --- TITLE PART ---
            if shape.has_text_frame:
                text = shape.text
                if "{{NAME}}" in text:
                    # Replace placeholder with user name
                    new_text = text.replace("{{NAME}}", chart_dict["NAME"].upper())
                    shape.text = new_text

                    # Apply formatting
                    for p in shape.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for run in p.runs:
                            run.font.size = Pt(14)
                            run.font.bold = True
                            run.font.shadow = True
                            run.font.color.rgb = RGBColor(0, 0, 0)  # black text
                            run.text = run.text.upper()

            # --- TABLE 1: Panchanga details ---
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        for key, val in chart_dict.items():
                            placeholder = "{{" + key + "}}"
                            if placeholder in text:
                                if hasattr(val, "strftime"):
                                    val = val.strftime("%Y-%m-%d %H:%M:%S")
                                text = text.replace(placeholder, str(val))
                        cell.text = text

                # --- Targeted override: set Karana/Yoga End by row label (no placeholders) ---
                try:
                    for row in table.rows:
                        cells = row.cells
                        if len(cells) >= 2:
                            label = cells[0].text.strip()
                            if label == "Karana End" and "KARANA_END" in chart_dict:
                                # Write the final string directly
                                cells[1].text = str(chart_dict["KARANA_END"])
                            elif label == "Yoga End" and "YOGA_END" in chart_dict:
                                cells[1].text = str(chart_dict["YOGA_END"])
                except Exception:
                    # Silent guard to avoid affecting other slides/tables
                    pass

            # --- TABLE 2: Rasi Chakra grid (limit to the 4x4 grid) ---
            if shape.has_table:
                table = shape.table

                # Guard: apply only to 4x4 Rasi Chakra table
                is_rasi_grid = (len(table.columns) == 4 and len(table.rows) == 4)
                if is_rasi_grid:
                    # print("Rasi Chakra augmentation applied to 4x4 table only.") 
                    for row in table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text in ZODIAC:
                                planets_here = []
                                for planet, lon in chart_dict["PLANETS"].items():
                                    sign = ZODIAC[int(lon // 30)]
                                    if sign == text:
                                        planets_here.append(planet_map[planet])
                                if planets_here:
                                    cell.text = text + "\n" + "\n".join(planets_here)
                                    # Center align appended lines
                                    if len(cell.text_frame.paragraphs) > 1:
                                        for p in cell.text_frame.paragraphs[1:]:
                                            p.alignment = PP_ALIGN.CENTER

            # --- TABLE 3: Planetary positions ---
            if shape.has_table:
                table = shape.table

                for row in table.rows:
                    for col_index, cell in enumerate(row.cells):  # enumerate to get column index
                        text = cell.text

                        for planet, suffix in planet_map.items():
                            # Column 2 → formatted D:M:S string
                            if f"{{{{LONG_{suffix}}}}}" in text:
                                val = chart_dict.get(f"LONG_{suffix}", "")
                                text = text.replace(f"{{{{LONG_{suffix}}}}}", val)

                            # Column 3 → raw longitude float
                            if f"{{{{ABS_{suffix}}}}}" in text:
                                val = chart_dict.get(f"ABS_{suffix}", "")
                                text = text.replace(f"{{{{ABS_{suffix}}}}}", val)

                            # Column 4 → zodiac sign abbreviation (computed from ABS longitude)
                            if f"{{{{SIGN_{suffix}}}}}" in text:
                                abs_val = float(chart_dict.get(f"ABS_{suffix}", 0))
                                sign_index = int(abs_val // 30)
                                val = ZODIAC[sign_index]
                                text = text.replace(f"{{{{SIGN_{suffix}}}}}", val)

                        # Write back if changed
                        if text != cell.text:
                            cell.text = text

                        # Center-align Zodiac_Sign column (4th column, index 3)
                        if col_index == 3:
                            for p in cell.text_frame.paragraphs:
                                p.alignment = PP_ALIGN.CENTER



    # --- Save with name + timestamp ---
    now_stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = chart_dict['NAME'].replace(" ", "_")
    output_file = f"{safe_name}_{now_stamp}.pptx"
    prs.save(output_file)
    print(f"✅ PPTX saved as {output_file}")

if __name__ == "__main__":
    # User input
    name = input("Enter Name: ")
    dob  = input("Enter Date of Birth (YYYY-MM-DD): ")
    tob  = input("Enter Time of Birth (HH:MM:SS): ")
    city = input("Enter Place of Birth: ")

    # Lookup city details
    # df = pd.read_excel("worldcities_with_timezone.xlsx")
    lat, lon, utc_offset, tz_string = lookup_city(city)

    tz = pytz.timezone(tz_string)

    # Convert DOB + TOB into datetime
    dob_obj = datetime.strptime(dob, "%Y-%m-%d").date()
    tob_obj = datetime.strptime(tob, "%H:%M:%S").time()
    dt_obj  = datetime.combine(dob_obj, tob_obj)
    

    # Compute Julian Day (local time)
    jd_local = swe.julday(
        dt_obj.year, dt_obj.month, dt_obj.day,
        dt_obj.hour + dt_obj.minute/60 + dt_obj.second/3600,
        swe.GREG_CAL
    )

    # Convert to UT
    jd_ut = jd_local - (utc_offset / 24.0)

    # Set sidereal mode (Lahiri ayanamsa)
    swe.set_sid_mode(swe.SIDM_LAHIRI, 0, 0)

    # Compute sidereal planetary longitudes (take [0] for longitude only)
    sun_lon     = swe.calc_ut(jd_ut, swe.SUN,     swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    moon_lon    = swe.calc_ut(jd_ut, swe.MOON,    swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    mars_lon    = swe.calc_ut(jd_ut, swe.MARS,    swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    mercury_lon = swe.calc_ut(jd_ut, swe.MERCURY, swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    jupiter_lon = swe.calc_ut(jd_ut, swe.JUPITER, swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    venus_lon   = swe.calc_ut(jd_ut, swe.VENUS,   swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]
    saturn_lon  = swe.calc_ut(jd_ut, swe.SATURN,  swe.FLG_SWIEPH | swe.FLG_SIDEREAL)[0][0]

    rahu_result = swe.calc_ut(jd_ut, swe.MEAN_NODE, swe.FLG_SWIEPH | swe.FLG_SIDEREAL)
    rahu_lon    = rahu_result[0][0]
    ketu_lon    = (rahu_lon + 180) % 360

    houses, ascmc = swe.houses(jd_ut, lat, lon, b'P')  # 'P' = Placidus
    asc_tropical = ascmc[0]
    ayan = swe.get_ayanamsa(jd_ut)
    asc_sidereal = (asc_tropical - ayan) % 360


    # Normalize to 0–360
    sidereal_positions = {
        "Sun": sun_lon % 360,
        "Moon": moon_lon % 360,
        "Mars": mars_lon % 360,
        "Mercury": mercury_lon % 360,
        "Jupiter": jupiter_lon % 360,
        "Venus": venus_lon % 360,
        "Saturn": saturn_lon % 360,
        "Rahu": rahu_lon % 360,
        "Ketu": ketu_lon % 360,
        "Ascendant": asc_sidereal   # ✅ included here
    }

    # Generate chart
    chart_dict = generate_birth_chart(
        name, dob, tob, city, lat, lon, utc_offset,
        jd_ut,
        sidereal_positions
    )